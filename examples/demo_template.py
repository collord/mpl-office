"""Template workflow demo.

Builds a small branded template in-place, then opens it and inserts a
matplotlib figure on a new content slide — the same workflow a real user
would follow with ``officer`` + ``rvg`` in R.

Run:

    .venv/Scripts/python examples/demo_template.py

Opens two files:
  - ``template.pptx`` — the branded shell
  - ``from_template.pptx`` — the template + a figure slide
"""
from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Pt

from mpl_office.pptx import Inches, fig_to_slide


HERE = Path(__file__).parent


def build_template(path: Path) -> None:
    """Write a 16:9 branded template with a title slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
    slide.shapes.title.text = "Quarterly Performance"

    # Subtitle placeholder
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Automated report - mpl-office demo"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(18)

    prs.save(path)
    print(f"wrote template ->{path}")


def build_report(template: Path, out: Path) -> None:
    """Open the template and append a figure slide."""
    prs = Presentation(template)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add a header text box on the content slide
    from pptx.util import Emu
    from pptx.util import Inches as PptxInches

    tb = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3),
        PptxInches(12.33), PptxInches(0.6),
    )
    tb.text_frame.text = "Revenue growth — Q1 through Q4"
    for para in tb.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # The figure itself
    fig, ax = plt.subplots(figsize=(10, 5.5))
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    revenue = [4.2, 5.1, 6.8, 7.5]
    forecast = [4.0, 5.3, 6.5, 8.0]
    x = np.arange(len(quarters))
    width = 0.35
    ax.bar(x - width / 2, revenue, width, label="actual", color="#2E86AB")
    ax.bar(x + width / 2, forecast, width, label="forecast", color="#A23B72")
    ax.set_xticks(x)
    ax.set_xticklabels(quarters)
    ax.set_ylabel("revenue ($M)")
    ax.legend()
    ax.grid(True, axis="y", alpha=0.3)

    fig_to_slide(
        fig, slide,
        left=Inches(1.5), top=Inches(1.2),
        width=Inches(10.33), height=Inches(5.5),
    )
    plt.close(fig)

    prs.save(out)
    print(f"wrote report   ->{out}")


def main():
    template = HERE / "template.pptx"
    report = HERE / "from_template.pptx"
    build_template(template)
    build_report(template, report)


if __name__ == "__main__":
    main()
