"""Verify `fig.savefig("out.pptx")` works via the mpl_office matplotlib backend."""
from __future__ import annotations

from pathlib import Path

import pytest

matplotlib = pytest.importorskip("matplotlib")
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402


def test_savefig_pptx(artifacts: Path):
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.plot([1, 2, 3], [1, 4, 9])
    ax.set_title("Backend test")
    out = artifacts / "backend_savefig.pptx"
    fig.savefig(out, backend="module://mpl_office.backend")
    plt.close(fig)

    assert out.exists() and out.stat().st_size > 1000
    prs = Presentation(out)
    assert len(prs.slides) == 1
    # Should contain at least one (grouped) shape
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) >= 1
