"""Build a small gallery .pptx with several matplotlib figures side-by-side.

Run:

    .venv/Scripts/python examples/demo.py

and open the resulting ``demo.pptx`` in PowerPoint for a visual check.
"""
from __future__ import annotations

import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation

from mpl_office.pptx import Inches, fig_to_slide


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: a simple line plot
    fig, ax = plt.subplots(figsize=(10, 5.5))
    x = np.linspace(0, 4 * np.pi, 200)
    ax.plot(x, np.sin(x), label="sin", linewidth=2)
    ax.plot(x, np.cos(x), label="cos", linewidth=2)
    ax.set_xlabel("angle (radians)")
    ax.set_ylabel("amplitude")
    ax.set_title("Trig functions")
    ax.legend()
    ax.grid(True, alpha=0.3)
    fig_to_slide(fig, _slide(prs), left=Inches(1.5), top=Inches(1),
                 width=Inches(10), height=Inches(5.5))
    plt.close(fig)

    # Slide 2: bar chart with edge colors
    fig, ax = plt.subplots(figsize=(8, 5))
    cats = ["Apples", "Oranges", "Bananas", "Cherries"]
    vals = [12, 19, 7, 15]
    ax.bar(cats, vals, color=["#E74C3C", "#F39C12", "#F1C40F", "#C0392B"],
           edgecolor="black", linewidth=1.5)
    ax.set_ylabel("count")
    ax.set_title("Fruit sales")
    fig_to_slide(fig, _slide(prs), left=Inches(2.5), top=Inches(1),
                 width=Inches(8), height=Inches(5))
    plt.close(fig)

    # Slide 3: scatter + annotation
    rng = np.random.default_rng(7)
    fig, ax = plt.subplots(figsize=(8, 5))
    xs = rng.uniform(0, 10, 100)
    ys = xs * 2 + rng.normal(0, 2, 100)
    ax.scatter(xs, ys, alpha=0.6, s=40)
    ax.set_xlabel("x")
    ax.set_ylabel("y")
    ax.set_title("Scatter with trend")
    ax.annotate("an outlier", xy=(xs.argmax(), ys[xs.argmax()]),
                xytext=(5, 25),
                arrowprops=dict(arrowstyle="->"))
    fig_to_slide(fig, _slide(prs), left=Inches(2.5), top=Inches(1),
                 width=Inches(8), height=Inches(5))
    plt.close(fig)

    # Slide 4: subplots
    fig, axes = plt.subplots(2, 2, figsize=(10, 6))
    for i, ax in enumerate(axes.flat):
        ax.plot(np.arange(10), np.arange(10) ** (i + 1))
        ax.set_title(f"power {i + 1}")
    fig.tight_layout()
    fig_to_slide(fig, _slide(prs), left=Inches(1.5), top=Inches(1),
                 width=Inches(10), height=Inches(6))
    plt.close(fig)

    out = Path(__file__).with_name("demo.pptx")
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
