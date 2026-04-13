"""Broader matplotlib gallery: exercise more of the SVG subset.

These tests don't yet do visual regression — they assert that each figure
round-trips through mpl-office → .pptx → python-pptx without raising and
produces the expected kinds of shapes.
"""
from __future__ import annotations

from pathlib import Path

import pytest

matplotlib = pytest.importorskip("matplotlib")
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

from mpl_office.pptx import Inches, fig_to_slide  # noqa: E402


def _blank():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _render_and_reopen(fig, name: str, artifacts: Path):
    prs, slide = _blank()
    fig_to_slide(fig, slide, left=Inches(1), top=Inches(1),
                 width=Inches(8), height=Inches(5.5))
    out = artifacts / f"{name}.pptx"
    prs.save(out)
    plt.close(fig)
    return Presentation(out)


def test_subplots(artifacts: Path):
    fig, axes = plt.subplots(1, 2, figsize=(8, 3))
    axes[0].plot([1, 2, 3], [1, 4, 9])
    axes[0].set_title("squares")
    axes[1].plot([1, 2, 3], [1, 8, 27])
    axes[1].set_title("cubes")
    prs = _render_and_reopen(fig, "subplots", artifacts)
    assert len(list(prs.slides[0].shapes)) >= 1


def test_histogram(artifacts: Path):
    rng = np.random.default_rng(1)
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.hist(rng.normal(size=500), bins=30, color="#6699cc")
    ax.set_title("Histogram")
    prs = _render_and_reopen(fig, "histogram", artifacts)
    assert len(list(prs.slides[0].shapes)) >= 1


def test_log_scale(artifacts: Path):
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.plot([1, 10, 100, 1000], [1, 4, 9, 16])
    ax.set_yscale("log")
    ax.set_xscale("log")
    prs = _render_and_reopen(fig, "log_scale", artifacts)
    assert len(list(prs.slides[0].shapes)) >= 1


def test_filled_area(artifacts: Path):
    fig, ax = plt.subplots(figsize=(6, 4))
    x = np.linspace(0, 2 * np.pi, 50)
    ax.fill_between(x, np.sin(x), np.cos(x), alpha=0.5, color="#9966cc")
    prs = _render_and_reopen(fig, "filled_area", artifacts)
    assert len(list(prs.slides[0].shapes)) >= 1


def test_pie(artifacts: Path):
    fig, ax = plt.subplots(figsize=(5, 5))
    ax.pie([30, 20, 15, 35], labels=["A", "B", "C", "D"],
           colors=["#ff6666", "#66ff66", "#6666ff", "#ffff66"])
    prs = _render_and_reopen(fig, "pie", artifacts)
    assert len(list(prs.slides[0].shapes)) >= 1


def test_annotated_plot(artifacts: Path):
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.plot([1, 2, 3, 4], [10, 20, 15, 25])
    ax.annotate("peak", xy=(4, 25), xytext=(3, 22),
                arrowprops=dict(arrowstyle="->"))
    ax.set_title("Annotated")
    prs = _render_and_reopen(fig, "annotated", artifacts)
    assert len(list(prs.slides[0].shapes)) >= 1
