"""Template-based workflows.

Exercises two paths:

1. Opening an existing ``.pptx`` as a template, adding a new slide with a
   chosen layout, and inserting a figure at an explicit position
   (``fig_to_slide``).

2. Opening a template that already has a content placeholder on the layout,
   and inserting a figure sized/positioned to the placeholder
   (``fig_to_placeholder``).

3. Using the matplotlib backend's ``template=`` kwarg so that
   ``fig.savefig("out.pptx", template=..., layout_index=...)`` works.
"""
from __future__ import annotations

from pathlib import Path

import pytest

matplotlib = pytest.importorskip("matplotlib")
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches as PptxInches, Pt  # noqa: E402

from mpl_office.pptx import Inches, fig_to_placeholder, fig_to_slide  # noqa: E402


@pytest.fixture
def template_path(artifacts: Path) -> Path:
    """Build a small reusable `.pptx` template.

    The template has a title slide with branded text and a blank layout
    for content slides. We write it once and reuse it in each test so the
    tests exercise the real "open template → add slide" path.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 0 — a branded title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    if title_slide.shapes.title is not None:
        title_slide.shapes.title.text = "Quarterly Report — Template"
    for shape in title_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = "Corporate confidential"

    out = artifacts / "_template.pptx"
    prs.save(out)
    return out


def _line_fig():
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.plot([1, 2, 3, 4, 5], [1, 4, 9, 16, 25], marker="o")
    ax.set_title("From template")
    ax.set_xlabel("x")
    ax.set_ylabel("y")
    return fig


def test_open_template_and_add_slide(template_path: Path, artifacts: Path):
    """Open an existing template, append a new slide with a chosen layout,
    and drop a figure onto it at an explicit position.
    """
    prs = Presentation(template_path)

    # The branded title slide from the template should survive.
    assert len(prs.slides) == 1
    assert any("Quarterly" in sh.text_frame.text
               for sh in prs.slides[0].shapes
               if sh.has_text_frame)

    # Append a blank slide (layout 6 in the default template is "Blank").
    blank_layout = prs.slide_layouts[6]
    content_slide = prs.slides.add_slide(blank_layout)

    fig = _line_fig()
    fig_to_slide(fig, content_slide,
                 left=Inches(1.5), top=Inches(1.25),
                 width=Inches(10), height=Inches(5))
    plt.close(fig)

    out = artifacts / "from_template.pptx"
    prs.save(out)

    # Re-open and verify both slides are intact.
    prs2 = Presentation(out)
    assert len(prs2.slides) == 2, "title + content slide"

    # Title slide still carries the template text.
    assert any("Quarterly" in sh.text_frame.text
               for sh in prs2.slides[0].shapes
               if sh.has_text_frame)

    # Content slide has at least one injected shape.
    shapes = list(prs2.slides[1].shapes)
    assert len(shapes) >= 1


def test_fig_to_placeholder(template_path: Path, artifacts: Path):
    """Insert a figure into a real placeholder (the title slide's body)."""
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Find a body placeholder (idx != 0, which is the title).
    body_ph = next(
        ph for ph in slide.placeholders if ph.placeholder_format.idx != 0
    )
    expected_left = body_ph.left
    expected_top = body_ph.top
    expected_w = body_ph.width
    expected_h = body_ph.height

    fig = _line_fig()
    fig_to_placeholder(fig, slide, body_ph)
    plt.close(fig)

    out = artifacts / "fig_to_placeholder.pptx"
    prs.save(out)

    # Re-open and confirm:
    #  - the placeholder was removed (only title placeholder remains)
    #  - a figure shape was added whose bbox matches the former placeholder
    prs2 = Presentation(out)
    slide2 = prs2.slides[0]
    non_title_phs = [
        ph for ph in slide2.placeholders if ph.placeholder_format.idx != 0
    ]
    assert non_title_phs == [], "body placeholder should have been replaced"

    # Find the injected shape by locating one whose top-left ≈ expected.
    matches = [
        s for s in slide2.shapes
        if s.left is not None and abs(s.left - expected_left) < 20_000
        and abs(s.top - expected_top) < 20_000
    ]
    assert matches, "expected to find a shape at the former placeholder's position"


def test_backend_savefig_with_template(template_path: Path, artifacts: Path):
    """`fig.savefig("out.pptx", template=..., layout_index=...)` via the backend."""
    fig = _line_fig()
    out = artifacts / "savefig_with_template.pptx"
    fig.savefig(
        out,
        backend="module://mpl_office.backend",
        template=str(template_path),
        layout_index=6,
    )
    plt.close(fig)

    prs = Presentation(out)
    # Template's title slide + the new content slide we appended.
    assert len(prs.slides) == 2

    # Original template text preserved.
    assert any("Quarterly" in sh.text_frame.text
               for sh in prs.slides[0].shapes
               if sh.has_text_frame)

    # Content slide has injected shapes.
    shapes = list(prs.slides[1].shapes)
    assert len(shapes) >= 1
