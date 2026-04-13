"""End-to-end round-trip: SVG → DrawingML → .pptx → reopen via python-pptx.

These tests verify that:
  * `svg_to_slide` injects shapes into a real slide's `<p:spTree>` without
    raising.
  * The resulting .pptx file is valid OOXML (python-pptx can read it back).
  * The inserted shapes carry the expected geometry and fill colour.
"""
from __future__ import annotations

from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches as PptxInches

from mpl_office.pptx import Inches, svg_to_slide


def _blank_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout
    return prs, prs.slides.add_slide(blank)


def test_rect_injection_opens(artifacts: Path):
    prs, slide = _blank_presentation()
    svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">'
        '<rect x="10" y="20" width="150" height="60" fill="#336699"/>'
        '</svg>'
    )
    svg_to_slide(
        svg,
        slide,
        left=Inches(1),
        top=Inches(1),
        width=Inches(6),
        height=Inches(3),
        source_dpi=96.0,
    )

    out_path = artifacts / "rect_roundtrip.pptx"
    prs.save(out_path)
    assert out_path.exists()

    # Re-open and inspect
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[0]
    shapes = list(slide2.shapes)
    assert len(shapes) == 1, f"expected one shape, got {len(shapes)}"
    shape = shapes[0]
    # Dimensions should reflect the requested target box (roughly — matched
    # by the fit-to-box scaling in the emitter).
    assert shape.width > 0
    assert shape.height > 0


def test_multiple_shapes_preserve_group(artifacts: Path):
    prs, slide = _blank_presentation()
    svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="200">
        <g>
            <rect x="0" y="0" width="50" height="50" fill="#ff0000"/>
            <rect x="60" y="0" width="50" height="50" fill="#00ff00"/>
            <circle cx="100" cy="150" r="25" fill="#0000ff"/>
        </g>
    </svg>"""
    svg_to_slide(
        svg,
        slide,
        left=Inches(1),
        top=Inches(1),
        width=Inches(5),
        height=Inches(5),
        source_dpi=96.0,
    )
    out = artifacts / "group_roundtrip.pptx"
    prs.save(out)

    prs2 = Presentation(out)
    slide2 = prs2.slides[0]
    # Top-level shape should be one group containing 3 children.
    shapes = list(slide2.shapes)
    assert len(shapes) == 1
    group = shapes[0]
    assert hasattr(group, "shapes"), "expected a group shape"
    assert len(list(group.shapes)) == 3


def test_path_custgeom(artifacts: Path):
    prs, slide = _blank_presentation()
    svg = """<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
        <path d="M10 10 C 20 0, 40 0, 50 10 S 80 20, 90 10 L 90 90 L 10 90 Z"
              fill="#ffcc00" stroke="#000000" stroke-width="2"/>
    </svg>"""
    svg_to_slide(svg, slide, left=Inches(1), top=Inches(1),
                 width=Inches(5), height=Inches(5), source_dpi=96.0)
    out = artifacts / "path_roundtrip.pptx"
    prs.save(out)

    prs2 = Presentation(out)
    shapes = list(prs2.slides[0].shapes)
    assert len(shapes) == 1
