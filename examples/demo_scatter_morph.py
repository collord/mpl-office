"""Build a two-slide .pptx wired for PowerPoint's Morph transition — a
Gapminder-style bubble scatter that animates between two years.

Each slide is the same scatter plot — X is fertility rate, Y is life
expectancy, bubble size scales with population — for the same five countries,
but at two different years (1960 and 1980). Because every bubble carries the
*same* shape name across both slides (via ``artist.set_gid("!!name")``),
PowerPoint's Morph transition pairs them up and tweens each bubble sliding to
its new position and growing/shrinking, instead of crossfading.

Two things have to be true for Morph to tween a bubble rather than crossfade
it: (1) the same shape *name* on both slides — done with ``set_gid("!!...")``,
as in ``demo_morph.py`` — and (2) the shape must sit at the slide's *top
level*. Morph matches the objects it can see in the Selection Pane; it pairs
*groups* as opaque units and won't descend into them. mpl-office renders the
figure as a nested group tree, so left alone every bubble crossfades AND the
static axes group fades out-and-back-in. ``_flatten_for_morph`` below dissolves
all groups so every shape — bubbles and chrome alike — is top-level and
individually matchable; the identical chrome then holds still while only the
bubbles travel.

Unlike ``demo_morph.py``, this script also writes the Morph transition into the
slide XML (see ``_inject_morph``), so no manual step is needed: just open
``demo_scatter_morph.pptx`` and advance through the slides to watch the bubbles
travel across the chart.

Run:

    .venv/bin/python examples/demo_scatter_morph.py        # macOS/Linux
    .venv/Scripts/python examples/demo_scatter_morph.py    # Windows
"""
from __future__ import annotations

import csv
from pathlib import Path

from lxml import etree

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation

from mpl_office.pptx import Inches, fig_to_slide

_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"

# (country code, display label, bubble colour)
COUNTRIES = [
    ("USA", "USA", "#3C6EB4"),
    ("FRA", "France", "#E74C3C"),
    ("IND", "India", "#F39C12"),
    ("BRA", "Brazil", "#27AE60"),
    ("CHN", "China", "#8E44AD"),
]

YEARS = ("1960","1970", "1980","1990","2000","2010")

# Fixed axes so a country's bubble lives at the same data coordinates whenever
# its values are the same — Morph reads pixel positions, so a stable frame is
# what makes the motion legible.
X_LIM = (0, 7.5)   # fertility rate, births per woman
Y_LIM = (30, 80)   # life expectancy, years

_DATA_DIR = Path(__file__).parent


def _load(filename):
    """Return ``{country_code: {year: float}}`` from a World Bank CSV."""
    out = {}
    with open(_DATA_DIR / filename, newline="", encoding="utf-8-sig") as fh:
        reader = csv.DictReader(fh)
        for row in reader:
            code = row["Country Code"]
            out[code] = {
                year: float(row[year]) for year in YEARS if row.get(year)
            }
    return out


def _bubble_area(population):
    """Map population to a scatter marker *area* in points^2.

    ``ax.scatter``'s ``s`` argument is an area in points^2, which is
    resolution- and aspect-independent — so bubbles stay round regardless of
    how the two axis scales differ. Area ∝ population is the honest encoding
    for a quantity; the divisor scales China (the largest) to a sensible
    on-slide size.
    """
    return population / 4.0e5


def _scatter_state(year, fertility, life_exp, population):
    """One figure for one year, with morph-matched shape names per country."""
    fig, ax = plt.subplots(figsize=(10, 5.5))

    for code, label, color in COUNTRIES:
        x = fertility[code][year]
        y = life_exp[code][year]
        pop = population[code][year]

        # One scatter call per country == one PathCollection carrying a clean
        # gid (it survives as <g id="!!bubble_CODE"> in the SVG). A *single*
        # scatter of all five points would collapse them into one collection
        # with one gid, which Morph can't pair per-country. Sizing via s=
        # (points^2 area) keeps bubbles round — a Circle patch with a data-unit
        # radius renders as an ellipse because the X and Y axis scales differ.
        ax.scatter([x], [y], s=_bubble_area(pop), color=color, alpha=0.7,
                   edgecolor="white", linewidth=0.5, zorder=3).set_gid(
                       f"!!bubble_{code}")

        ax.annotate(label, xy=(x, y), xytext=(0, 6),
                    textcoords="offset points", ha="center",
                    fontsize=9, color="#222222").set_gid(f"!!label_{code}")

    ax.set_xlim(*X_LIM)
    ax.set_ylim(*Y_LIM)
    ax.set_xlabel("Fertility rate (births per woman)")
    ax.set_ylabel("Life expectancy (years)")
    ax.set_title(f"Fertility vs. life expectancy — {year}")
    ax.grid(True, alpha=0.3)
    return fig


def _flatten_for_morph(slide):
    """Dissolve every group so Morph sees a flat list of top-level shapes.

    PowerPoint's Morph transition only matches the objects it can see at the
    top of the Selection Pane — it pairs *groups* as opaque units and never
    descends inside them. mpl-office faithfully renders the figure as a nested
    tree (``figure_1`` > ``axes_1`` > ticks/grid/bubbles), and that nesting
    breaks Morph two ways:

    * the bubbles, buried in the groups, never get paired point-to-point, so
      they crossfade instead of travelling; and
    * the *static* chrome (axes, grid, tick labels — pixel-identical on both
      slides) is trapped in a group that Morph can't reconcile, so the whole
      group fades out and back in even though nothing in it changed.

    Dissolving all groups fixes both: every shape becomes a first-class,
    individually-matchable object. mpl-office already bakes absolute
    coordinates into each child's ``<a:off>`` (group ``chOff`` == ``off``), so
    flattening doesn't move anything.

    One ordering subtlety: the axes background ``patch`` is an opaque white
    rectangle. After flattening it must stay *below* the bubbles in z-order
    (document order), or it paints over them. We re-append the ``!!``-named
    bubbles/labels last so they sit on top.
    """
    sp_tree = slide.shapes._spTree
    grp = f"{{{_PML}}}grpSp"
    sp = f"{{{_PML}}}sp"
    pic = f"{{{_PML}}}pic"
    cnvpr = f"{{{_PML}}}cNvPr"
    leaf_tags = (sp, pic, grp)

    # 1. Lift every leaf out of its group (innermost first), preserving order,
    #    then drop the emptied group shells.
    moved = True
    while moved:
        moved = False
        for group in list(sp_tree.iter(grp)):
            for child in list(group):
                if child.tag in leaf_tags:
                    group.remove(child)
                    group.addprevious(child)
                    moved = True
            if all(c.tag not in leaf_tags for c in group):
                group.getparent().remove(group)

    # 2. Restack the morph-tagged shapes on top of the (now top-level) chrome.
    for shape in [
        s for s in sp_tree.findall(sp)
        if (props := s.find(f".//{cnvpr}")) is not None
        and props.get("name", "").startswith("!!")
    ]:
        sp_tree.remove(shape)
        sp_tree.append(shape)


# The Morph transition element, wrapped in <mc:AlternateContent>. This is the
# form modern PowerPoint authors itself: a <mc:Choice Requires="p159"> holding
# the real <p159:morph>, plus a plain-fade <mc:Fallback> for renderers that
# don't understand it. A bare <p14:morph> (no wrapper) is silently downgraded
# to Fade on open in current PowerPoint — the wrapper is what makes the
# transition reliably register as Morph. ``{p}`` is filled with the
# presentationml namespace at injection time.
_MORPH_TRANSITION = (
    '<mc:AlternateContent'
    ' xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"'
    ' xmlns:p="{p}">'
    '<mc:Choice'
    ' xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main"'
    ' Requires="p159">'
    '<p:transition spd="slow" p14:dur="1500"'
    ' xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">'
    '<p159:morph option="byObject"/>'
    '</p:transition>'
    '</mc:Choice>'
    '<mc:Fallback>'
    '<p:transition xmlns:p="{p}" spd="slow"><p:fade/></p:transition>'
    '</mc:Fallback>'
    '</mc:AlternateContent>'
)


def _inject_morph(slide):
    """Set the slide's incoming transition to Morph, in place.

    Saves the manual "select slide, Transitions > Morph" step. The transition
    belongs to the slide it advances *into*, so this is called on slides 2..N.
    It's inserted right after ``<p:cSld>`` (and ``<p:clrMapOvr>`` if present),
    which is where the schema requires the transition element to sit.
    """
    sld = slide._element
    transition = etree.fromstring(_MORPH_TRANSITION.format(p=_PML))
    c_sld = sld.find(f"{{{_PML}}}cSld")
    clr_map = sld.find(f"{{{_PML}}}clrMapOvr")
    (clr_map if clr_map is not None else c_sld).addnext(transition)


def build():
    fertility = _load("fertility_rate.csv")
    life_exp = _load("life_expectancy.csv")
    population = _load("country_population.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for index, year in enumerate(YEARS):
        fig = _scatter_state(year, fertility, life_exp, population)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fig_to_slide(fig, slide, left=Inches(1.5), top=Inches(1),
                     width=Inches(10), height=Inches(5.5))
        _flatten_for_morph(slide)
        # The first slide has no incoming transition; every later slide morphs
        # from the one before it.
        if index > 0:
            _inject_morph(slide)
        plt.close(fig)

    out = Path(__file__).with_name("demo_scatter_morph.pptx")
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
