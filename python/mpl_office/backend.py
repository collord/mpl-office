"""matplotlib backend — write ``.pptx`` (and eventually ``.docx``) files
directly from ``fig.savefig(...)``.

Register by calling::

    import matplotlib
    matplotlib.use("module://mpl_office.backend")

or via the ``backend`` kwarg::

    fig.savefig("out.pptx", backend="module://mpl_office.backend")
"""
from __future__ import annotations

from io import StringIO
from typing import Any, Optional

from matplotlib.backend_bases import FigureManagerBase, _Backend
from matplotlib.backends.backend_svg import FigureCanvasSVG

from .pptx import Inches, svg_to_slide


class FigureCanvasOffice(FigureCanvasSVG):
    """Canvas subclass adding ``.pptx`` / ``.docx`` output formats."""

    # Tell matplotlib which formats we support so ``savefig("x.pptx")`` wires
    # us up instead of complaining about an unknown extension.
    filetypes = {**FigureCanvasSVG.filetypes, "pptx": "PowerPoint", "docx": "Word"}

    def get_default_filetype(self) -> str:  # pragma: no cover
        return "pptx"

    def print_pptx(
        self,
        filename,
        *,
        template: Optional[str] = None,
        layout_index: int = 6,
        placeholder_index: Optional[int] = None,
        **kwargs: Any,
    ) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches as PptxInches
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError(
                "python-pptx is required for .pptx output; "
                "install with `pip install python-pptx` or the `pptx` extra."
            ) from exc

        buf = StringIO()
        self.print_svg(buf)
        svg = buf.getvalue()

        # Figure dimensions in EMU.
        w_in, h_in = self.figure.get_size_inches()
        width_emu = Inches(w_in)
        height_emu = Inches(h_in)

        if template:
            prs = Presentation(template)
            try:
                layout = prs.slide_layouts[layout_index]
            except IndexError:
                layout = prs.slide_layouts[-1]
            slide = prs.slides.add_slide(layout)
        else:
            prs = Presentation()
            prs.slide_width = width_emu
            prs.slide_height = height_emu
            blank_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank_layout)

        if placeholder_index is not None and template:
            ph = slide.placeholders[placeholder_index]
            svg_to_slide(
                svg,
                slide,
                left=ph.left,
                top=ph.top,
                width=ph.width,
                height=ph.height,
                source_dpi=72.0,
            )
        else:
            svg_to_slide(
                svg,
                slide,
                left=0,
                top=0,
                width=width_emu,
                height=height_emu,
                source_dpi=72.0,
            )

        prs.save(filename)

    def print_docx(self, filename, **kwargs: Any) -> None:  # pragma: no cover
        raise NotImplementedError(".docx output is not yet implemented")


@_Backend.export
class _BackendOffice(_Backend):
    FigureCanvas = FigureCanvasOffice
    FigureManager = FigureManagerBase
