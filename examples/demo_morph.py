"""Build a two-slide .pptx wired for PowerPoint's Morph transition.

The trick: graphics elements that represent *the same data point* on both
slides carry the *same* shape name. We do that by calling
``artist.set_gid("!!name")`` on each artist — the ``!!`` prefix survives
matplotlib's SVG backend as an ``id`` attribute, and ``mpl-office`` maps it
to ``<p:cNvPr name="!!name">`` in the DrawingML output. PowerPoint's Morph
transition matches shapes across consecutive slides *by name*, so a bar
tagged ``!!bar_Q1`` on slide 1 and ``!!bar_Q1`` on slide 2 is understood as
the same object in two states — Morph tweens its height instead of
crossfading.

This script only wires up the *shape names* that Morph matches on; it does
not set the transition itself. After running, open ``demo_morph.pptx``,
select slide 2, and apply Transitions > Morph in PowerPoint. Then advancing
from slide 1 to slide 2 animates each bar growing/shrinking and the marker
sliding across.

(Writing the Morph transition into the slide XML directly was tried and
removed — the injected ``<p159:morph>`` element did not reliably activate
the transition when opened in PowerPoint, so the manual step is the
supported path for now.)

Run:

    .venv/bin/python examples/demo_morph.py        # macOS/Linux
    .venv/Scripts/python examples/demo_morph.py    # Windows
"""
from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Circle

from pptx import Presentation

from mpl_office.pptx import Inches, fig_to_slide

CATEGORIES = ["Q1", "Q2", "Q3", "Q4"]

# Two states of the same series. Matching index == matching data point.
STATE_A = [4.2, 5.1, 6.8, 7.5]
STATE_B = [6.1, 3.4, 8.9, 5.0]

# A single marker that travels to a new (x, y) between the two slides.
MARKER_A = (0, 2.0)
MARKER_B = (3, 8.0)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bar_state(values, marker_xy):
    """One figure for one state, with morph-matched shape names.

    Every bar gets gid ``!!bar_<category>`` and the travelling marker gets
    ``!!marker`` — identical names across both states, so Morph pairs them.
    """
    fig, ax = plt.subplots(figsize=(10, 5.5))

    bars = ax.bar(CATEGORIES, values, color="#2E86AB", edgecolor="black",
                  linewidth=1.0)
    for cat, bar in zip(CATEGORIES, bars):
        bar.set_gid(f"!!bar_{cat}")

    # A travelling marker. Drawn as a Circle patch — a single <path> that
    # carries !!marker cleanly. (ax.plot's "o" marker also works now that
    # set_gid propagates through <use>-backed artists, but a patch keeps the
    # demo's one-shape-one-name mapping unambiguous.)
    marker = Circle(marker_xy, radius=0.18, color="#E74C3C", zorder=5)
    ax.add_patch(marker)
    marker.set_gid("!!marker")

    ax.set_ylim(0, 10)
    ax.set_ylabel("revenue ($M)")
    ax.set_title("Quarterly revenue")
    return fig


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for values, marker in ((STATE_A, MARKER_A), (STATE_B, MARKER_B)):
        fig = _bar_state(values, marker)
        fig_to_slide(fig, _blank_slide(prs), left=Inches(1.5), top=Inches(1),
                     width=Inches(10), height=Inches(5.5))
        plt.close(fig)

    out = Path(__file__).with_name("demo_morph.pptx")
    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
