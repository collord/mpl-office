"""Raster `<image>` support.

Exercises:

1. The low-level data-URI round trip (SVG with an inline base64 PNG).
2. A matplotlib ``imshow`` heatmap — the common real-world case, where
   the SVG backend emits an embedded PNG that our converter has to
   decode, re-embed as an OOXML image part, and reference via a real
   relationship id.
3. A figure with both vector shapes (axes, tick labels) **and** an
   imshow raster, to confirm the two paths coexist on the same slide.
"""
from __future__ import annotations

import base64
import zipfile
from io import BytesIO
from pathlib import Path

import pytest

matplotlib = pytest.importorskip("matplotlib")
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

import mpl_office  # noqa: E402
from mpl_office.pptx import Inches, fig_to_slide, svg_to_slide  # noqa: E402


# 1×1 red PNG — smallest useful payload for pure-unit data URI tests.
RED_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)
RED_PNG_BYTES = base64.b64decode(RED_PNG_B64)


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _list_image_parts(pptx_path: Path) -> list[str]:
    """Return paths of every media file embedded in the .pptx."""
    with zipfile.ZipFile(pptx_path) as z:
        return [n for n in z.namelist() if n.startswith("ppt/media/")]


def _read_image_part(pptx_path: Path, part_name: str) -> bytes:
    with zipfile.ZipFile(pptx_path) as z:
        return z.read(part_name)


def test_convert_with_images_extracts_data_uri():
    """The low-level API returns image bytes alongside XML."""
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg"
             xmlns:xlink="http://www.w3.org/1999/xlink"
             width="100" height="100">
        <image x="10" y="10" width="80" height="80"
               xlink:href="data:image/png;base64,{RED_PNG_B64}"/>
    </svg>"""
    xml, images = mpl_office.convert_svg_to_drawingml_with_images(svg)

    assert len(images) == 1
    sentinel, data, fmt = images[0]
    assert sentinel == "__mpl_office_img_0__"
    assert fmt == "png"
    assert data == RED_PNG_BYTES
    # The XML references the sentinel, not a real relationship id.
    assert f'r:embed="{sentinel}"' in xml
    assert "<p:pic>" in xml


def test_convert_legacy_api_drops_images_silently():
    """The string-only API is backwards-compatible: same XML, no images."""
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg"
             xmlns:xlink="http://www.w3.org/1999/xlink"
             width="100" height="100">
        <image x="0" y="0" width="100" height="100"
               xlink:href="data:image/png;base64,{RED_PNG_B64}"/>
    </svg>"""
    xml = mpl_office.convert_svg_to_drawingml(svg)
    # The picture shape + sentinel are still in the XML — the caller just
    # doesn't get the image bytes, which means if they inject without
    # rewriting sentinels the picture will have a dangling r:embed.
    # That's the documented contract of the legacy function.
    assert "<p:pic>" in xml


def test_raw_svg_with_image_roundtrip(artifacts: Path):
    """Raw SVG with an inline PNG survives a full save/reopen cycle."""
    prs, slide = _blank_slide()
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg"
             xmlns:xlink="http://www.w3.org/1999/xlink"
             width="200" height="200">
        <rect x="0" y="0" width="200" height="200" fill="#eeeeee"/>
        <image x="20" y="20" width="160" height="160"
               xlink:href="data:image/png;base64,{RED_PNG_B64}"/>
    </svg>"""
    svg_to_slide(
        svg, slide,
        left=Inches(1), top=Inches(1),
        width=Inches(5), height=Inches(5),
        source_dpi=96.0,
    )
    out = artifacts / "raw_image.pptx"
    prs.save(out)

    # The .pptx should contain exactly one media file, and it should be
    # byte-identical to the source PNG.
    media = _list_image_parts(out)
    assert len(media) == 1
    assert _read_image_part(out, media[0]) == RED_PNG_BYTES

    # And reopening via python-pptx should not raise.
    prs2 = Presentation(out)
    assert len(list(prs2.slides[0].shapes)) >= 1


@pytest.fixture
def imshow_fig():
    fig, ax = plt.subplots(figsize=(5, 5))
    rng = np.random.default_rng(0)
    ax.imshow(rng.random((20, 20)), cmap="viridis", interpolation="nearest")
    ax.set_title("Random heatmap")
    ax.set_xlabel("x")
    ax.set_ylabel("y")
    yield fig
    plt.close(fig)


def test_imshow_embeds_png(imshow_fig, artifacts: Path):
    """An imshow figure round-trips with the bitmap preserved in the .pptx."""
    prs, slide = _blank_slide()
    fig_to_slide(
        imshow_fig, slide,
        left=Inches(1), top=Inches(1),
        width=Inches(5), height=Inches(5),
    )
    out = artifacts / "imshow.pptx"
    prs.save(out)

    media = _list_image_parts(out)
    assert media, "expected imshow to produce at least one media part"
    # Each media file should be a recognizable PNG.
    for m in media:
        blob = _read_image_part(out, m)
        assert blob.startswith(b"\x89PNG\r\n\x1a\n"), f"{m} is not a PNG"

    # Re-open and confirm the slide has both vector shapes and a picture.
    prs2 = Presentation(out)
    slide2 = prs2.slides[0]

    def walk(shapes):
        pics = 0
        others = 0
        for s in shapes:
            if hasattr(s, "shapes"):
                p, o = walk(s.shapes)
                pics += p
                others += o
            else:
                # shape_type 13 is PICTURE in python-pptx
                if str(s.shape_type).startswith("PICTURE"):
                    pics += 1
                else:
                    others += 1
        return pics, others

    pics, others = walk(slide2.shapes)
    assert pics >= 1, "expected at least one picture shape from imshow"
    assert others >= 1, "expected axis/title shapes alongside the picture"


def test_duplicate_images_deduplicated(artifacts: Path):
    """Two SVG <image> refs to the same data URI share one image part."""
    prs, slide = _blank_slide()
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg"
             xmlns:xlink="http://www.w3.org/1999/xlink"
             width="200" height="200">
        <image x="0"  y="0"  width="80" height="80"
               xlink:href="data:image/png;base64,{RED_PNG_B64}"/>
        <image x="100" y="100" width="80" height="80"
               xlink:href="data:image/png;base64,{RED_PNG_B64}"/>
    </svg>"""
    svg_to_slide(
        svg, slide,
        left=Inches(1), top=Inches(1),
        width=Inches(5), height=Inches(5),
        source_dpi=96.0,
    )
    out = artifacts / "dup_image.pptx"
    prs.save(out)

    # python-pptx's ImagePart cache deduplicates by content hash, so even
    # though we pushed two <p:pic> shapes we should have exactly one
    # media file in the package.
    media = _list_image_parts(out)
    assert len(media) == 1, f"expected one deduplicated media file, got {media}"
