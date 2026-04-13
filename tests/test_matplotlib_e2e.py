"""End-to-end: render a real matplotlib figure, convert to .pptx, reopen.

This is the smoke test that actually exercises the full target workflow.
If matplotlib's SVG backend produces something our converter chokes on,
this test catches it.
"""
from __future__ import annotations

from pathlib import Path

import pytest

matplotlib = pytest.importorskip("matplotlib")
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

from mpl_office.pptx import Inches, fig_to_slide  # noqa: E402


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


@pytest.fixture
def lineplot_fig():
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.plot([1, 2, 3, 4, 5], [1, 4, 9, 16, 25], marker="o", label="squares")
    ax.plot([1, 2, 3, 4, 5], [1, 8, 27, 64, 125], marker="s", label="cubes")
    ax.set_xlabel("x")
    ax.set_ylabel("y")
    ax.set_title("Polynomial growth")
    ax.legend()
    yield fig
    plt.close(fig)


@pytest.fixture
def bar_fig():
    fig, ax = plt.subplots(figsize=(5, 3))
    ax.bar(["a", "b", "c", "d"], [3, 7, 5, 2], color=["#ff0000", "#00ff00", "#0000ff", "#ffaa00"])
    ax.set_title("Bar")
    yield fig
    plt.close(fig)


@pytest.fixture
def scatter_fig():
    import numpy as np
    rng = np.random.default_rng(42)
    fig, ax = plt.subplots(figsize=(5, 5))
    ax.scatter(rng.normal(size=50), rng.normal(size=50), c=rng.uniform(size=50), cmap="viridis")
    yield fig
    plt.close(fig)


def test_lineplot_roundtrip(lineplot_fig, artifacts: Path):
    prs, slide = _blank_slide()
    fig_to_slide(lineplot_fig, slide, left=Inches(1), top=Inches(1),
                 width=Inches(8), height=Inches(5))
    out = artifacts / "lineplot.pptx"
    prs.save(out)

    prs2 = Presentation(out)
    slide2 = prs2.slides[0]
    shapes = list(slide2.shapes)
    assert len(shapes) >= 1, "expected at least one injected shape"


def test_bar_roundtrip(bar_fig, artifacts: Path):
    prs, slide = _blank_slide()
    fig_to_slide(bar_fig, slide, left=Inches(1), top=Inches(1),
                 width=Inches(6), height=Inches(4))
    out = artifacts / "bar.pptx"
    prs.save(out)

    prs2 = Presentation(out)
    shapes = list(prs2.slides[0].shapes)
    assert len(shapes) >= 1


def test_scatter_roundtrip(scatter_fig, artifacts: Path):
    prs, slide = _blank_slide()
    fig_to_slide(scatter_fig, slide, left=Inches(1), top=Inches(1),
                 width=Inches(5), height=Inches(5))
    out = artifacts / "scatter.pptx"
    prs.save(out)

    prs2 = Presentation(out)
    shapes = list(prs2.slides[0].shapes)
    assert len(shapes) >= 1
